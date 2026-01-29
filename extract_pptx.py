
import os
from pptx import Presentation

# List of files provided by user
files = [
    "Chap01(1) (3).pptx",
    "Chap02 (3).pptx",
    "Chap03 (3).pptx",
    "Chap04 (4).pptx",
    "Chap05(1) (3).pptx",
    "Chap06(1) (2).pptx",
    "Chap08 (3).pptx",
    "Chap12(1) (2).pptx",
    "Chap13(1) (1).pptx",
    "Chap14(1) (1).pptx"
]

base_path = r"c:\Dev\assaf-professional-portfolio"

print("Parsing PPTX files to extract Case Study topics...\n")

for filename in files:
    path = os.path.join(base_path, filename)
    if not os.path.exists(path):
        print(f"Skipping {filename} (Not found)")
        continue

    try:
        prs = Presentation(path)
        title = ""
        subtitle = ""
        topics = []
        
        # Try to get title from first slide
        if len(prs.slides) > 0:
            slide1 = prs.slides[0]
            if slide1.shapes.title:
                title = slide1.shapes.title.text.strip()
            # Try to find subtitle/content in other shapes of slide 1
            for shape in slide1.shapes:
                if shape.has_text_frame and shape.text != title:
                     if len(shape.text.strip()) > 0:
                        subtitle = shape.text.strip().replace("\n", " ")

        # Scan other slides for outlines or key topics
        # Heuristic: Look for slides with bullets
        content_text = []
        for i, slide in enumerate(prs.slides[1:5]): # Check first few slides for an agenda
             if slide.shapes.title:
                 header = slide.shapes.title.text.strip()
                 content_text.append(header)
        
        print(f"--- FILE: {filename} ---")
        print(f"TITLE: {title}")
        print(f"SUBTITLE: {subtitle[:100]}...") # truncate
        print(f"KEY TOPICS (Headers): {content_text}")
        print("\n")
        
    except Exception as e:
        print(f"Error parsing {filename}: {e}\n")
